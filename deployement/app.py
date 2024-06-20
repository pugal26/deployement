from flask import Flask, request, redirect, render_template, send_file
import os
import shutil
import cv2
import time
from skimage.metrics import structural_similarity as ssim
from pptx import Presentation
from pptx.util import Inches
import numpy as np
from keras.applications import VGG16
from keras.applications.vgg16 import preprocess_input
from keras.preprocessing import image as keras_image
from werkzeug.utils import secure_filename

app = Flask(__name__)

# Function to create a directory if it doesn't exist
def create_folder(directory):
    if not os.path.exists(directory):
        os.makedirs(directory)

# Function to extract frames from a video
def extract_frames(video_path, output_folder, frame_interval=30):
    vidcap = cv2.VideoCapture(video_path)
    success, image = vidcap.read()
    frame_count = 0
    saved_frame_count = 0
    start_time = time.time()
    while success:
        if frame_count % frame_interval == 0:
            cv2.imwrite(os.path.join(output_folder, f"frame{saved_frame_count:04d}.jpg"), image)
            saved_frame_count += 1
        success, image = vidcap.read()
        frame_count += 1
    end_time = time.time()
    duration = end_time - start_time
    vidcap.release()
    return saved_frame_count, duration

# Function to calculate image similarity using SSIM
def calculate_image_similarity(image1, image2):
    gray1 = cv2.cvtColor(image1, cv2.COLOR_BGR2GRAY)
    gray2 = cv2.cvtColor(image2, cv2.COLOR_BGR2GRAY)
    score, _ = ssim(gray1, gray2, full=True)
    return score

# Function to create PowerPoint presentation from images in a folder
def create_ppt(folder_path, output_ppt):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    image_files = [os.path.join(folder_path, filename) for filename in os.listdir(folder_path) if filename.endswith(('.png', '.jpg', '.jpeg'))]
    image_files.sort()
    for image_file in image_files:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        left = top = 0
        slide.shapes.add_picture(image_file, left, top, width=prs.slide_width, height=prs.slide_height)
    prs.save(output_ppt)

# Function to calculate VGG16 feature vector for an image
def calculate_image_feature(image_path):
    img = keras_image.load_img(image_path, target_size=(224, 224))
    img_array = keras_image.img_to_array(img)
    img_array = np.expand_dims(img_array, axis=0)
    img_array = preprocess_input(img_array)
    feature = vgg_model.predict(img_array)
    return feature.flatten()

# Function to delete all files in a folder
def delete_all_files_in_folder(folder_path):
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        if os.path.isfile(file_path) or os.path.islink(file_path):
            os.unlink(file_path)
        elif os.path.isdir(file_path):
            shutil.rmtree(file_path)

# Load pre-trained VGG16 model
vgg_model = VGG16(weights='imagenet', include_top=False, input_shape=(224, 224, 3))

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return redirect(request.url)
    file = request.files['file']
    if file.filename == '':
        return redirect(request.url)
    if file:
        video_filename = secure_filename(file.filename)
        video_path = os.path.join('uploads', video_filename)
        file.save(video_path)

        output_folder = os.path.join('processed', os.path.splitext(video_filename)[0])
        create_folder(output_folder)

        frame_count, extraction_duration = extract_frames(video_path, output_folder, frame_interval=30)
        print(f"Processed video '{os.path.basename(video_path)}': {frame_count} frames extracted in {extraction_duration} seconds")

        output_folder_similar = os.path.join('similar', os.path.splitext(video_filename)[0])
        create_folder(output_folder_similar)

        image_files = [os.path.join(output_folder, filename) for filename in os.listdir(output_folder) if filename.endswith(('.png', '.jpg', '.jpeg'))]
        image_files.sort()

        image_features = {}
        for image_file in image_files:
            image_features[image_file] = calculate_image_feature(image_file)

        threshold = 0.95
        for i in range(len(image_files) - 1):
            current_image_file = image_files[i]
            current_image_feature = image_features[current_image_file]
            
            next_image_file = image_files[i + 1]
            next_image_feature = image_features[next_image_file]
            
            similarity_score = np.dot(current_image_feature, next_image_feature) / (np.linalg.norm(current_image_feature) * np.linalg.norm(next_image_feature))
            
            if similarity_score > threshold:
                print(f"Images {os.path.basename(current_image_file)} and {os.path.basename(next_image_file)} are highly similar (Cosine Similarity: {similarity_score:.2f}). Moving {os.path.basename(current_image_file)}.")
                shutil.move(current_image_file, os.path.join(output_folder_similar, os.path.basename(current_image_file)))

        pptx_path = os.path.join('pptx', os.path.splitext(video_filename)[0] + '.pptx')
        create_folder('pptx')
        create_ppt(output_folder, pptx_path)

        delete_all_files_in_folder(output_folder)
        delete_all_files_in_folder(output_folder_similar)

        return send_file(pptx_path, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
